from pptx import Presentation

from dedoc.data_structures.paragraph_metadata import ParagraphMetadata
from dedoc.data_structures.table import Table
from dedoc.data_structures.table_metadata import TableMetadata
from dedoc.data_structures.unstructured_document import UnstructuredDocument
from dedoc.readers.base_reader import BaseReader
from dedoc.data_structures.line_with_meta import LineWithMeta
from dedoc.extensions import recognized_extensions, recognized_mimes
from dedoc.readers.utils.hierarch_level_extractor import HierarchyLevelExtractor

from typing import Optional, Tuple


class PptxReader(BaseReader):
    def __init__(self):
        self.hierarchy_level_extractor = HierarchyLevelExtractor()

    def can_read(self, path: str, mime: str, extension: str, document_type: str) -> bool:
        return extension in recognized_extensions.pptx_like_format or mime in recognized_mimes.pptx_like_format

    def read(self,
             path: str,
             document_type: Optional[str] = None,
             parameters: Optional[dict] = None) -> Tuple[UnstructuredDocument, bool]:
        prs = Presentation(path)
        lines, tables = [], []

        for page_id, slide in enumerate(prs.slides, start=1):
            for paragraph_id, shape in enumerate(slide.shapes, start=1):

                if shape.has_text_frame:
                    metadata = ParagraphMetadata(paragraph_type="raw_text",
                                                 predicted_classes=None,
                                                 page_id=page_id,
                                                 line_id=paragraph_id)
                    lines.append(LineWithMeta(line=shape.text, hierarchy_level=None, metadata=metadata, annotations=[]))

                if shape.has_table:
                    cells = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    metadata = TableMetadata(page_id=page_id)
                    tables.append(Table(cells=cells, metadata=metadata))

        lines = self.hierarchy_level_extractor.get_hierarchy_level(lines)
        return UnstructuredDocument(lines=lines, tables=tables), True
